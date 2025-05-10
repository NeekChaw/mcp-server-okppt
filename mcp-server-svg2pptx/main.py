from mcp.server.fastmcp import FastMCP, Context
from pptx.util import Inches, Pt, Cm, Emu
from typing import Optional, Union, List, Tuple
import os
import datetime
import traceback
import re
from svg_module import insert_svg_to_pptx, to_emu, create_svg_file, get_pptx_slide_count, save_svg_code_to_file

# 创建MCP服务器实例
mcp = FastMCP(name="main")

# 路径辅助函数
def get_base_dir():
    """获取基础目录（服务器目录的父目录）"""
    current_dir = os.path.dirname(os.path.abspath(__file__))
    return os.path.dirname(current_dir)

def get_tmp_dir():
    """获取临时文件目录，如果不存在则创建"""
    tmp_dir = os.path.join(get_base_dir(), "tmp")
    os.makedirs(tmp_dir, exist_ok=True)
    return tmp_dir

def get_output_dir():
    """获取输出文件目录，如果不存在则创建"""
    output_dir = os.path.join(get_base_dir(), "output")
    os.makedirs(output_dir, exist_ok=True)
    return output_dir

def cleanup_filename(filename: str) -> str:
    """
    清理文件名，移除所有旧的时间戳和操作类型标记
    
    Args:
        filename: 要清理的文件名（不含路径和扩展名）
        
    Returns:
        清理后的基本文件名
    """
    # 移除类似 _svg_20240101_120000, _deleted_20240529_153045 等操作标记和时间戳
    # 模式: _ + 操作名 + _ + 8位日期 + _ + 6位时间
    pattern = r'_(svg|deleted|inserted|output)_\d{8}_\d{6}'
    cleaned = re.sub(pattern, '', filename)
    
    # 防止文件名连续处理后残留多余的下划线
    cleaned = re.sub(r'_{2,}', '_', cleaned)
    
    # 移除末尾的下划线(如果有)
    cleaned = cleaned.rstrip('_')
    
    return cleaned

def get_default_output_path(file_type="pptx", base_name=None, op_type=None):
    """
    获取默认输出文件路径
    
    Args:
        file_type: 文件类型（扩展名）
        base_name: 基本文件名，如果为None则使用时间戳
        op_type: 操作类型，用于在文件名中添加标记
    
    Returns:
        默认输出文件路径
    """
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    
    if base_name is None:
        base_name = f"output_{timestamp}"
    else:
        # 清理基本文件名
        base_name = cleanup_filename(base_name)
        
        # 添加操作类型和时间戳
        if op_type:
            base_name = f"{base_name}_{op_type}_{timestamp}"
        else:
            base_name = f"{base_name}_{timestamp}"
    
    return os.path.join(get_output_dir(), f"{base_name}.{file_type}")

# 主要的SVG插入工具
@mcp.tool()
def insert_svg(
    pptx_path: str,# 空字符串表示自动创建，否则使用绝对路径
    svg_path: List[str],# 数组，绝对路径
    slide_number: int = 1,
    x_inches: float = 0,
    y_inches: float = 0,
    width_inches: float = 16,
    height_inches: float = 9,
    output_path: str = "",# 空字符串表示自动创建，否则使用绝对路径
    create_if_not_exists: bool = True
) -> str:
    """
    将SVG图像插入到PPTX文件的指定位置。(如果需要替换已有的幻灯片，请组合使用`delete_slide`和`insert_blank_slide`功能)
    如果未提供PPTX路径，将自动创建一个临时文件，位于服务器同级目录的tmp目录。
    如果未提供输出路径，将使用标准输出目录，位于服务器同级目录的output目录。
    如果未提供坐标，默认对齐幻灯片左上角。
    如果未提供宽度和高度，默认覆盖整个幻灯片（16:9）。

    支持批量处理：
    - 如果svg_path是单个字符串数组，则将SVG添加到slide_number指定的页面
    - 如果svg_path是列表，则从slide_number开始顺序添加每个SVG，即第一个SVG添加到
      slide_number页，第二个添加到slide_number+1页，依此类推

    Args:
        pptx_path: PPTX文件路径，如果未提供则自动创建一个临时文件，最好使用英文路径
        svg_path: SVG文件路径或SVG文件路径列表，最好使用英文路径
        slide_number: 起始幻灯片编号（从1开始）
        x_inches: X坐标（英寸），如果未指定则默认为0
        y_inches: Y坐标（英寸），如果未指定则默认为0
        width_inches: 宽度（英寸），如果未指定则使用幻灯片宽度
        height_inches: 高度（英寸），如果未指定则根据宽度计算或使用幻灯片高度
        output_path: 输出文件路径，如果未指定则使用标准输出目录
        create_if_not_exists: 如果为True且PPTX文件不存在，将自动创建一个新文件
        
    Returns:
        操作结果消息，包含详细的错误信息（如果有）
    """
    # 收集错误信息
    error_messages = []
    result_messages = []

    # 如果未提供pptx_path，使用默认输出目录创建一个
    if not pptx_path or pptx_path.strip() == "":
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        pptx_path = os.path.join(get_output_dir(), f"presentation_{timestamp}.pptx")
        print(f"未提供PPTX路径，将使用默认路径: {pptx_path}")

    # 处理输出路径
    if not output_path:
        # 从原始文件名生成输出文件名
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        base_name = cleanup_filename(base_name)
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(get_output_dir(), f"{base_name}_svg_{timestamp}.pptx")

    if not os.path.isabs(pptx_path):
        pptx_path = os.path.abspath(pptx_path)
    
    # 确保PPTX文件的父目录存在
    pptx_dir = os.path.dirname(pptx_path)
    if not os.path.exists(pptx_dir):
        try:
            os.makedirs(pptx_dir, exist_ok=True)
            print(f"已创建PPTX目录: {pptx_dir}")
            error_messages.append(f"已创建PPTX目录: {pptx_dir}")
        except Exception as e:
            error_msg = f"创建PPTX目录 {pptx_dir} 时出错: {e}"
            error_messages.append(error_msg)
            return error_msg
    
    # 将英寸转换为Inches对象
    x = Inches(x_inches) if x_inches is not None else None
    y = Inches(y_inches) if y_inches is not None else None
    width = Inches(width_inches) if width_inches is not None else None
    height = Inches(height_inches) if height_inches is not None else None
    
    # 如果提供了输出路径且是相对路径，转换为绝对路径
    if output_path and not os.path.isabs(output_path):
        output_path = os.path.abspath(output_path)
    
    # 如果提供了输出路径，确保其父目录存在
    if output_path:
        output_dir = os.path.dirname(output_path)
        if not os.path.exists(output_dir):
            try:
                os.makedirs(output_dir, exist_ok=True)
                print(f"已创建输出目录: {output_dir}")
                error_messages.append(f"已创建输出目录: {output_dir}")
            except Exception as e:
                error_msg = f"创建输出目录 {output_dir} 时出错: {e}"
                error_messages.append(error_msg)
                return error_msg
    
    # 检查svg_path的类型并分别处理
    if isinstance(svg_path, str):
        # 单个SVG文件处理
        return process_single_svg(
            pptx_path, svg_path, slide_number, x, y, width, height, 
            output_path, create_if_not_exists
        )
    elif isinstance(svg_path, list):
        # 批量处理SVG文件列表
        success_count = 0
        total_count = len(svg_path)
        
        if total_count == 0:
            return "错误：SVG文件列表为空"
        
        # 创建中间文件路径基础
        temp_base = os.path.join(get_tmp_dir(), f"svg_batch_{datetime.datetime.now().strftime('%Y%m%d_%H%M%S')}")
        os.makedirs(os.path.dirname(temp_base), exist_ok=True)
        
        # 当前输入文件路径
        current_input = pptx_path
        
        for i, current_svg in enumerate(svg_path):
            current_slide = slide_number + i
            
            # 处理每个SVG文件
            if i < total_count - 1:
                # 对于非最后一个文件，创建临时输出路径
                temp_output = f"{temp_base}_step_{i}.pptx"
                
                result = process_single_svg(
                    current_input,
                    current_svg, 
                    current_slide, 
                    x, y, width, height, 
                    temp_output, 
                    create_if_not_exists
                )
                
                # 下一次迭代的输入文件是本次的输出文件
                current_input = temp_output
            else:
                # 最后一个SVG使用最终输出路径
                final_output = output_path if output_path else pptx_path
                
                result = process_single_svg(
                    current_input,
                    current_svg, 
                    current_slide, 
                    x, y, width, height, 
                    final_output, 
                    create_if_not_exists
                )
            
            # 检查处理结果
            if "成功" in result:
                success_count += 1
                result_messages.append(f"第{i+1}个SVG({current_svg})：成功添加到第{current_slide}页")
            else:
                result_messages.append(f"第{i+1}个SVG({current_svg})：添加失败 - {result}")
        
        # 清理临时文件
        for i in range(total_count - 1):
            temp_file = f"{temp_base}_step_{i}.pptx"
            if os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                except Exception as e:
                    print(f"清理临时文件 {temp_file} 时出错: {e}")
        
        # 返回总体结果
        result_path = output_path or pptx_path
        summary = f"批量处理完成：共{total_count}个SVG文件，成功{success_count}个，失败{total_count-success_count}个"
        details = "\n".join(result_messages)
        return f"{summary}\n输出文件：{result_path}\n\n详细结果：\n{details}"
    else:
        return f"错误：svg_path类型无效，必须是字符串或字符串列表，当前类型: {type(svg_path)}"

def process_single_svg(
    pptx_path: str,
    svg_path: str,
    slide_number: int,
    x: Optional[Union[Inches, Pt, Cm, Emu, float]],
    y: Optional[Union[Inches, Pt, Cm, Emu, float]],
    width: Optional[Union[Inches, Pt, Cm, Emu, float]],
    height: Optional[Union[Inches, Pt, Cm, Emu, float]],
    output_path: Optional[str],
    create_if_not_exists: bool
) -> str:
    """处理单个SVG文件的辅助函数"""
    # 检查SVG文件是否存在，如果是相对路径则转换为绝对路径
    if not os.path.isabs(svg_path):
        svg_path = os.path.abspath(svg_path)
    
    # 确保SVG文件的父目录存在
    svg_dir = os.path.dirname(svg_path)
    if not os.path.exists(svg_dir):
        try:
            os.makedirs(svg_dir, exist_ok=True)
            print(f"已创建SVG目录: {svg_dir}")
        except Exception as e:
            return f"创建SVG目录 {svg_dir} 时出错: {e}"
        
    # 如果SVG文件不存在且create_if_not_exists为True，则创建一个简单的SVG文件
    if not os.path.exists(svg_path) and create_if_not_exists:
        svg_created = create_svg_file(svg_path)
        if not svg_created:
            return f"错误：无法创建SVG文件 {svg_path}"
    elif not os.path.exists(svg_path):
        return f"错误：SVG文件 {svg_path} 不存在"
    
    # 确保输出路径存在，如果未指定则使用标准输出目录
    if not output_path:
        # 从原始文件名生成输出文件名
        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        # 清理文件名
        base_name = cleanup_filename(base_name)
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = os.path.join(get_output_dir(), f"{base_name}_svg_{timestamp}.pptx")
    
    try:
        # 调用改进后的函数，它现在返回一个元组 (成功标志, 错误消息)
        result = insert_svg_to_pptx(
            pptx_path=pptx_path,
            svg_path=svg_path,
            slide_number=slide_number,
            x=x,
            y=y,
            width=width,
            height=height,
            output_path=output_path,
            create_if_not_exists=create_if_not_exists
        )
        
        # 检查返回值类型
        if isinstance(result, tuple) and len(result) == 2:
            success, error_details = result
        else:
            # 向后兼容
            success = result
            error_details = ""
        
        if success:
            result_path = output_path or pptx_path
            was_created = not os.path.exists(pptx_path) and create_if_not_exists
            creation_msg = "（已自动创建PPTX文件）" if was_created else ""
            return f"成功将SVG文件 {svg_path} 插入到 {result_path} 的第 {slide_number} 张幻灯片 {creation_msg}"
        else:
            # 返回详细的错误信息
            return f"插入SVG到PPTX文件失败，详细错误信息：\n{error_details}"
    except Exception as e:
        # 收集异常堆栈
        error_trace = traceback.format_exc()
        return f"插入SVG时发生错误: {str(e)}\n\n详细堆栈跟踪：\n{error_trace}"

@mcp.tool()
def list_files(directory: str = ".", file_type: Optional[str] = None) -> str:
    """
    列出目录中的文件，可选按文件类型过滤。
    
    Args:
        directory: 要列出文件的目录路径
        file_type: 文件类型过滤，可以是 "svg" 或 "pptx"
        
    Returns:
        文件列表（每行一个文件）
    """
    import os
    
    if not os.path.exists(directory):
        return f"错误：目录 {directory} 不存在"
    
    files = os.listdir(directory)
    
    if file_type:
        file_type = file_type.lower()
        extensions = {
            "svg": [".svg"],
            "pptx": [".pptx", ".ppt"]
        }
        
        if file_type in extensions:
            filtered_files = []
            for file in files:
                if any(file.lower().endswith(ext) for ext in extensions[file_type]):
                    filtered_files.append(file)
            files = filtered_files
        else:
            files = [f for f in files if f.lower().endswith(f".{file_type}")]
    
    if not files:
        return f"未找到{'任何' if not file_type else f'{file_type}'} 文件"
    
    return "\n".join(files)

@mcp.tool()
def get_file_info(file_path: str) -> str:
    """
    获取文件信息，如存在状态、大小等。
    
    Args:
        file_path: 要查询的文件路径
        
    Returns:
        文件信息
    """
    import os
    
    if not os.path.exists(file_path):
        return f"文件 {file_path} 不存在"
    
    if os.path.isdir(file_path):
        return f"{file_path} 是一个目录"
    
    size_bytes = os.path.getsize(file_path)
    size_kb = size_bytes / 1024
    size_mb = size_kb / 1024
    
    if size_mb >= 1:
        size_str = f"{size_mb:.2f} MB"
    else:
        size_str = f"{size_kb:.2f} KB"
    
    modified_time = os.path.getmtime(file_path)
    from datetime import datetime
    modified_str = datetime.fromtimestamp(modified_time).strftime("%Y-%m-%d %H:%M:%S")
    
    # 获取文件扩展名
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    file_type = None
    if ext == ".svg":
        file_type = "SVG图像"
    elif ext in [".pptx", ".ppt"]:
        file_type = "PowerPoint演示文稿"
    else:
        file_type = f"{ext[1:]} 文件" if ext else "未知类型文件"
    
    return f"文件: {file_path}\n类型: {file_type}\n大小: {size_str}\n修改时间: {modified_str}"

# 添加一个将SVG转换为PNG的工具
@mcp.tool()
def convert_svg_to_png(
    svg_path: str,
    output_path: Optional[str] = None
) -> str:
    """
    将SVG文件转换为PNG图像。
    
    Args:
        svg_path: SVG文件路径
        output_path: 输出PNG文件路径，如果未指定则使用相同文件名但扩展名为.png
        
    Returns:
        操作结果消息
    """
    from reportlab.graphics import renderPM
    from svglib.svglib import svg2rlg
    import os
    
    if not os.path.exists(svg_path):
        return f"错误：SVG文件 {svg_path} 不存在"
    
    if not output_path:
        # 获取不带扩展名的文件名，然后添加.png扩展名
        base_name = os.path.splitext(svg_path)[0]
        output_path = f"{base_name}.png"
    
    try:
        drawing = svg2rlg(svg_path)
        if drawing is None:
            return f"错误：无法读取SVG文件 {svg_path}"
        
        renderPM.drawToFile(drawing, output_path, fmt="PNG")
        return f"成功将SVG文件 {svg_path} 转换为PNG文件 {output_path}\n宽度: {drawing.width}px\n高度: {drawing.height}px"
    except Exception as e:
        return f"转换SVG到PNG时发生错误: {str(e)}"

@mcp.tool()
def get_pptx_info(pptx_path: str) -> str:
    """
    获取PPTX文件的基本信息，包括幻灯片数量。
    
    Args:
        pptx_path: PPTX文件路径
        
    Returns:
        包含文件信息和幻灯片数量的字符串
    """
    import os
    
    # 确保路径存在
    if not os.path.isabs(pptx_path):
        pptx_path = os.path.abspath(pptx_path)
    
    # 先获取基本文件信息
    if not os.path.exists(pptx_path):
        return f"错误：文件 {pptx_path} 不存在"
    
    size_bytes = os.path.getsize(pptx_path)
    size_kb = size_bytes / 1024
    size_mb = size_kb / 1024
    
    if size_mb >= 1:
        size_str = f"{size_mb:.2f} MB"
    else:
        size_str = f"{size_kb:.2f} KB"
    
    modified_time = os.path.getmtime(pptx_path)
    from datetime import datetime
    modified_str = datetime.fromtimestamp(modified_time).strftime("%Y-%m-%d %H:%M:%S")
    
    # 获取幻灯片数量
    slide_count, error = get_pptx_slide_count(pptx_path)
    
    if error:
        slide_info = f"获取幻灯片数量失败：{error}"
    else:
        slide_info = f"幻灯片数量：{slide_count}张"
    
    return f"PPT文件: {pptx_path}\n大小: {size_str}\n修改时间: {modified_str}\n{slide_info}"

@mcp.tool()
def save_svg_code(
    svg_code: str
) -> str:
    """
    将SVG代码保存为SVG文件并返回保存的绝对路径。
    !!!注意：特殊字符如"&"需要转义为"&amp;"
    Args:
        svg_code: SVG代码内容
        
    Returns:
        操作结果消息，包含保存的文件路径或错误信息
    """
    try:
        # 调用svg_module中的函数保存SVG代码
        success, file_path, error_message = save_svg_code_to_file(
            svg_code=svg_code,
            output_path="",
            create_dirs=True
        )
        
        if success:
            return f"成功保存SVG代码到文件: {file_path}"
        else:
            return f"保存SVG代码到文件失败: {error_message}"
    except Exception as e:
        error_trace = traceback.format_exc()
        return f"保存SVG代码到文件时发生错误: {str(e)}\n\n详细堆栈跟踪：\n{error_trace}"

@mcp.tool()
def delete_slide(
    pptx_path: str,
    slide_number: int,
    output_path: Optional[str] = None
) -> str:
    """
    从PPTX文件中删除指定编号的幻灯片。

    !!!注意：

    在使用SVG替换PPT幻灯片内容时，我们发现了一些关键点，以下是正确替换PPT内容的方法总结：

    ### 正确替换PPT内容的方法

    1. **完全替换法**（最可靠）：
    - 删除需要替换的幻灯片（使用`delete_slide`功能）
    - 在同一位置插入空白幻灯片（使用`insert_blank_slide`功能）
    - 将新的SVG内容插入到空白幻灯片（使用`insert_svg`功能）

    2. **新文件法**（适合多页修改）：
    - 创建全新的PPT文件
    - 将所有需要的SVG（包括已修改的）按顺序插入到新文件中
    - 这样可以避免在旧文件上操作导致的混淆和叠加问题

    3. **注意事项**：
    - 直接对现有幻灯片插入SVG会导致新内容叠加在原内容上，而非替换
    - 文件名可能会随着多次操作变得过长，影响可读性
    - 批量插入SVG时，`svg_path`参数必须是数组形式，即使只有一个文件
    - 操作后应检查输出文件以确认修改是否成功

    ### 推荐工作流

    1. 先保存修改后的SVG内容到文件
    2. 创建一个全新的PPT文件
    3. 按顺序一次性插入所有SVG（包括已修改和未修改的）
    4. 使用简洁直观的文件名

    这种方法避免了多步骤操作导致的文件混乱，也能确保每张幻灯片都是干净的、不包含叠加内容的。

    Args:
        pptx_path: PPTX文件路径
        slide_number: 要删除的幻灯片编号（从1开始）
        output_path: 输出文件路径，如果未指定则使用标准输出目录
        
    Returns:
        操作结果消息
    """
    try:
        # 确保路径是绝对路径
        if not os.path.isabs(pptx_path):
            pptx_path = os.path.abspath(pptx_path)
            
        # 检查文件是否存在
        if not os.path.exists(pptx_path):
            return f"错误：PPTX文件 {pptx_path} 不存在"
            
        # 处理输出路径，如果未指定则使用标准输出目录
        if not output_path:
            # 从原始文件名生成输出文件名
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            # 清理文件名
            base_name = cleanup_filename(base_name)
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(get_output_dir(), f"{base_name}_deleted_{timestamp}.pptx")
            
        if output_path and not os.path.isabs(output_path):
            output_path = os.path.abspath(output_path)
            
        # 如果提供了输出路径，确保其父目录存在
        if output_path:
            output_dir = os.path.dirname(output_path)
            if not os.path.exists(output_dir):
                try:
                    os.makedirs(output_dir, exist_ok=True)
                except Exception as e:
                    return f"创建输出目录 {output_dir} 时出错: {e}"
        
        # 使用python-pptx加载演示文稿
        from pptx import Presentation
        prs = Presentation(pptx_path)
        
        # 检查幻灯片编号范围
        if not 1 <= slide_number <= len(prs.slides):
            return f"错误：幻灯片编号 {slide_number} 超出范围 [1, {len(prs.slides)}]"
        
        # 计算索引（转换为从0开始）
        slide_index = slide_number - 1
        
        # 使用用户提供的方法删除幻灯片
        slides = list(prs.slides._sldIdLst)
        prs.slides._sldIdLst.remove(slides[slide_index])
        
        # 保存文件
        save_path = output_path
        prs.save(save_path)
        
        return f"成功从 {pptx_path} 中删除第 {slide_number} 张幻灯片，结果已保存到 {save_path}"
        
    except Exception as e:
        error_trace = traceback.format_exc()
        return f"删除幻灯片时发生错误: {str(e)}\n\n详细堆栈跟踪：\n{error_trace}"

@mcp.tool()
def insert_blank_slide(
    pptx_path: str,
    slide_number: int,
    layout_index: int = 6,  # 默认使用空白布局
    output_path: Optional[str] = None,
    create_if_not_exists: bool = True
) -> str:
    """
    在PPTX文件的指定位置插入一个空白幻灯片。

    !!!注意：

    在使用SVG替换PPT幻灯片内容时，我们发现了一些关键点，以下是正确替换PPT内容的方法总结：

    ### 正确替换PPT内容的方法

    1. **完全替换法**（最可靠）：
    - 删除需要替换的幻灯片（使用`delete_slide`功能）
    - 在同一位置插入空白幻灯片（使用`insert_blank_slide`功能）
    - 将新的SVG内容插入到空白幻灯片（使用`insert_svg`功能）

    2. **新文件法**（适合多页修改）：
    - 创建全新的PPT文件
    - 将所有需要的SVG（包括已修改的）按顺序插入到新文件中
    - 这样可以避免在旧文件上操作导致的混淆和叠加问题

    3. **注意事项**：
    - 直接对现有幻灯片插入SVG会导致新内容叠加在原内容上，而非替换
    - 文件名可能会随着多次操作变得过长，影响可读性
    - 批量插入SVG时，`svg_path`参数必须是数组形式，即使只有一个文件
    - 操作后应检查输出文件以确认修改是否成功

    ### 推荐工作流

    1. 先保存修改后的SVG内容到文件
    2. 创建一个全新的PPT文件
    3. 按顺序一次性插入所有SVG（包括已修改和未修改的）
    4. 使用简洁直观的文件名

    这种方法避免了多步骤操作导致的文件混乱，也能确保每张幻灯片都是干净的、不包含叠加内容的。

    Args:
        pptx_path: PPTX文件路径
        slide_number: 要插入幻灯片的位置编号（从1开始）
        layout_index: 幻灯片布局索引，默认为6（空白布局）
        output_path: 输出文件路径，如果未指定则使用标准输出目录
        create_if_not_exists: 如果为True且PPTX文件不存在，将自动创建一个新文件
        
    Returns:
        操作结果消息
    """
    try:
        # 如果未提供pptx_path，使用默认输出目录创建一个
        if not pptx_path or pptx_path.strip() == "":
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            pptx_path = os.path.join(get_output_dir(), f"presentation_{timestamp}.pptx")
            print(f"未提供PPTX路径，将使用默认路径: {pptx_path}")
            
        # 确保路径是绝对路径
        if not os.path.isabs(pptx_path):
            pptx_path = os.path.abspath(pptx_path)
            
        # 处理输出路径，如果未指定则使用标准输出目录
        if not output_path:
            # 从原始文件名生成输出文件名
            base_name = os.path.splitext(os.path.basename(pptx_path))[0]
            # 清理文件名
            base_name = cleanup_filename(base_name)
            timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = os.path.join(get_output_dir(), f"{base_name}_inserted_{timestamp}.pptx")
            
        if output_path and not os.path.isabs(output_path):
            output_path = os.path.abspath(output_path)
            
        # 如果提供了输出路径，确保其父目录存在
        if output_path:
            output_dir = os.path.dirname(output_path)
            if not os.path.exists(output_dir):
                try:
                    os.makedirs(output_dir, exist_ok=True)
                except Exception as e:
                    return f"创建输出目录 {output_dir} 时出错: {e}"
        
        # 检查文件是否存在
        file_exists = os.path.exists(pptx_path)
        if not file_exists and not create_if_not_exists:
            return f"错误：PPTX文件 {pptx_path} 不存在，且未启用自动创建"
            
        # 使用python-pptx加载或创建演示文稿
        from pptx import Presentation
        prs = Presentation(pptx_path) if file_exists else Presentation()
        
        # 如果是新创建的演示文稿，设置为16:9尺寸
        if not file_exists:
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
        
        # 检查布局索引是否有效
        if layout_index >= len(prs.slide_layouts):
            return f"错误：无效的布局索引 {layout_index}，可用范围 [0, {len(prs.slide_layouts)-1}]"
        
        # 检查幻灯片编号范围
        slides_count = len(prs.slides)
        if not 1 <= slide_number <= slides_count + 1:
            return f"错误：幻灯片位置 {slide_number} 超出范围 [1, {slides_count + 1}]"
        
        # 计算索引（转换为从0开始）
        slide_index = slide_number - 1
        
        # 在末尾添加新幻灯片
        new_slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        # 如果不是添加到末尾，需要移动幻灯片
        if slide_index < slides_count:
            # 获取幻灯片列表
            slides = list(prs.slides._sldIdLst)
            # 将最后一张幻灯片（刚添加的）移动到目标位置
            last_slide = slides[-1]
            # 从列表中移除最后一张幻灯片
            prs.slides._sldIdLst.remove(last_slide)
            # 在目标位置插入幻灯片
            prs.slides._sldIdLst.insert(slide_index, last_slide)
        
        # 保存文件
        save_path = output_path
        prs.save(save_path)
        
        # 构建返回消息
        action = "添加" if file_exists else "创建并添加"
        return f"成功在 {pptx_path} 中{action}第 {slide_number} 张幻灯片，结果已保存到 {save_path}"
        
    except Exception as e:
        error_trace = traceback.format_exc()
        return f"插入幻灯片时发生错误: {str(e)}\n\n详细堆栈跟踪：\n{error_trace}"


# 启动服务器
if __name__ == "__main__":
    # 确保必要的目录存在
    tmp_dir = get_tmp_dir()
    output_dir = get_output_dir()
    print(f"临时文件目录: {tmp_dir}")
    print(f"输出文件目录: {output_dir}")
    
    # 启动MCP服务器
    print("\n启动MCP服务器...")
    mcp.run() 