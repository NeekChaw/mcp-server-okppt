from mcp.server.fastmcp import FastMCP, Context
from pptx.util import Inches, Pt, Cm, Emu
from typing import Optional, Union, List
import os
from svg_module import insert_svg_to_pptx, to_emu

# 创建MCP服务器实例
mcp = FastMCP(name="SVG-to-PPTX")

# 主要的SVG插入工具
@mcp.tool()
def insert_svg(
    pptx_path: str,
    svg_path: str,
    slide_number: int = 1,
    x_inches: Optional[float] = None,
    y_inches: Optional[float] = None,
    width_inches: Optional[float] = None,
    height_inches: Optional[float] = None,
    output_path: Optional[str] = None,
    create_if_not_exists: bool = True
) -> str:
    """
    将SVG图像插入到PPTX文件的指定位置。
    如果未提供PPTX路径，将自动创建一个临时文件。
    如果未提供输出路径，将覆盖原始文件。
    
    Args:
        pptx_path: PPTX文件路径，如果未提供则自动创建一个临时文件
        svg_path: SVG文件路径
        slide_number: 目标幻灯片编号（从1开始）
        x_inches: X坐标（英寸），如果未指定则默认为0
        y_inches: Y坐标（英寸），如果未指定则默认为0
        width_inches: 宽度（英寸），如果未指定则使用幻灯片宽度
        height_inches: 高度（英寸），如果未指定则根据宽度计算或使用幻灯片高度
        output_path: 输出文件路径，如果未指定则覆盖原始文件
        create_if_not_exists: 如果为True且PPTX文件不存在，将自动创建一个新文件
        
    Returns:
        操作结果消息
    """
    import datetime
    
    # 检查pptx_path是否为空，如果为空则创建默认路径
    using_default_path = False
    if not pptx_path or pptx_path.strip() == "":
        # 创建一个基于时间戳的默认文件名
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        pptx_path = f"presentation_{timestamp}.pptx"
        print(f"未提供PPTX路径，将使用默认路径: {pptx_path}")
        using_default_path = True
    
    # 将英寸转换为Inches对象
    x = Inches(x_inches) if x_inches is not None else None
    y = Inches(y_inches) if y_inches is not None else None
    width = Inches(width_inches) if width_inches is not None else None
    height = Inches(height_inches) if height_inches is not None else None
    
    # 检查SVG文件是否存在
    if not os.path.exists(svg_path):
        return f"错误：SVG文件 {svg_path} 不存在"
    
    try:
        # 调用原始函数
        success = insert_svg_to_pptx(
            pptx_path=pptx_path,
            svg_path=svg_path,
            slide_number=slide_number,
            x=x,
            y=y,
            width=width,
            height=height,
            output_path=output_path,
            create_if_not_exists=create_if_not_exists
        )
        
        if success:
            result_path = output_path or pptx_path
            was_created = not os.path.exists(pptx_path) and create_if_not_exists
            creation_msg = "（已自动创建PPTX文件）" if was_created else ""
            default_path_msg = "（使用自动生成的默认路径）" if using_default_path else ""
            return f"成功将SVG文件 {svg_path} 插入到 {result_path} 的第 {slide_number} 张幻灯片 {creation_msg}{default_path_msg}"
        else:
            return f"插入SVG到PPTX文件失败，详细错误请查看日志"
    except Exception as e:
        return f"插入SVG时发生错误: {str(e)}"

@mcp.tool()
def list_files(directory: str = ".", file_type: Optional[str] = None) -> str:
    """
    列出目录中的文件，可选按文件类型过滤。
    
    Args:
        directory: 要列出文件的目录路径
        file_type: 文件类型过滤，可以是 "svg" 或 "pptx"
        
    Returns:
        文件列表（每行一个文件）
    """
    import os
    
    if not os.path.exists(directory):
        return f"错误：目录 {directory} 不存在"
    
    files = os.listdir(directory)
    
    if file_type:
        file_type = file_type.lower()
        extensions = {
            "svg": [".svg"],
            "pptx": [".pptx", ".ppt"]
        }
        
        if file_type in extensions:
            filtered_files = []
            for file in files:
                if any(file.lower().endswith(ext) for ext in extensions[file_type]):
                    filtered_files.append(file)
            files = filtered_files
        else:
            files = [f for f in files if f.lower().endswith(f".{file_type}")]
    
    if not files:
        return f"未找到{'任何' if not file_type else f'{file_type}'} 文件"
    
    return "\n".join(files)

@mcp.tool()
def get_file_info(file_path: str) -> str:
    """
    获取文件信息，如存在状态、大小等。
    
    Args:
        file_path: 要查询的文件路径
        
    Returns:
        文件信息
    """
    import os
    
    if not os.path.exists(file_path):
        return f"文件 {file_path} 不存在"
    
    if os.path.isdir(file_path):
        return f"{file_path} 是一个目录"
    
    size_bytes = os.path.getsize(file_path)
    size_kb = size_bytes / 1024
    size_mb = size_kb / 1024
    
    if size_mb >= 1:
        size_str = f"{size_mb:.2f} MB"
    else:
        size_str = f"{size_kb:.2f} KB"
    
    modified_time = os.path.getmtime(file_path)
    from datetime import datetime
    modified_str = datetime.fromtimestamp(modified_time).strftime("%Y-%m-%d %H:%M:%S")
    
    # 获取文件扩展名
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    file_type = None
    if ext == ".svg":
        file_type = "SVG图像"
    elif ext in [".pptx", ".ppt"]:
        file_type = "PowerPoint演示文稿"
    else:
        file_type = f"{ext[1:]} 文件" if ext else "未知类型文件"
    
    return f"文件: {file_path}\n类型: {file_type}\n大小: {size_str}\n修改时间: {modified_str}"

# 添加一个将SVG转换为PNG的工具
@mcp.tool()
def convert_svg_to_png(
    svg_path: str,
    output_path: Optional[str] = None
) -> str:
    """
    将SVG文件转换为PNG图像。
    
    Args:
        svg_path: SVG文件路径
        output_path: 输出PNG文件路径，如果未指定则使用相同文件名但扩展名为.png
        
    Returns:
        操作结果消息
    """
    from reportlab.graphics import renderPM
    from svglib.svglib import svg2rlg
    import os
    
    if not os.path.exists(svg_path):
        return f"错误：SVG文件 {svg_path} 不存在"
    
    if not output_path:
        # 获取不带扩展名的文件名，然后添加.png扩展名
        base_name = os.path.splitext(svg_path)[0]
        output_path = f"{base_name}.png"
    
    try:
        drawing = svg2rlg(svg_path)
        if drawing is None:
            return f"错误：无法读取SVG文件 {svg_path}"
        
        renderPM.drawToFile(drawing, output_path, fmt="PNG")
        return f"成功将SVG文件 {svg_path} 转换为PNG文件 {output_path}\n宽度: {drawing.width}px\n高度: {drawing.height}px"
    except Exception as e:
        return f"转换SVG到PNG时发生错误: {str(e)}"

# 添加一个批量处理SVG文件的工具
@mcp.tool()
def batch_insert_svgs(
    pptx_path: str,
    svg_dir: str,
    slide_number: int = 1,
    output_path: Optional[str] = None,
    create_if_not_exists: bool = True
) -> str:
    """
    将目录中的所有SVG文件批量插入到PPTX文件中（每个SVG插入一张新幻灯片）。
    
    Args:
        pptx_path: PPTX文件路径，如果未提供则自动创建一个临时文件
        svg_dir: 包含SVG文件的目录
        slide_number: 开始插入的幻灯片编号
        output_path: 输出文件路径，如果未指定则覆盖原始文件
        create_if_not_exists: 如果为True且PPTX文件不存在，将自动创建一个新文件
        
    Returns:
        操作结果消息
    """
    import os
    from pptx import Presentation
    import datetime
    
    # 检查pptx_path是否为空，如果为空则创建默认路径
    if not pptx_path or pptx_path.strip() == "":
        # 创建一个基于时间戳的默认文件名
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        pptx_path = f"presentation_{timestamp}.pptx"
        print(f"未提供PPTX路径，将使用默认路径: {pptx_path}")
    
    if not os.path.exists(pptx_path) and create_if_not_exists:
        try:
            # 创建新的PPTX文件
            prs = Presentation()
            # 设置为16:9尺寸
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            # 添加一张空白幻灯片
            blank_slide_layout = prs.slide_layouts[6]  # 6是空白幻灯片
            slide = prs.slides.add_slide(blank_slide_layout)
            prs.save(pptx_path)
            print(f"自动创建PPTX文件: {pptx_path}")
        except Exception as e:
            return f"创建PPTX文件时出错: {str(e)}"
    elif not os.path.exists(pptx_path):
        return f"错误：PPTX文件 {pptx_path} 不存在"
    
    if not os.path.exists(svg_dir) or not os.path.isdir(svg_dir):
        return f"错误：目录 {svg_dir} 不存在或不是一个目录"
    
    # 获取目录中的所有SVG文件
    svg_files = [f for f in os.listdir(svg_dir) if f.lower().endswith('.svg')]
    
    if not svg_files:
        return f"错误：目录 {svg_dir} 中没有SVG文件"
    
    # 排序SVG文件以确保顺序一致
    svg_files.sort()
    
    # 准备输出路径
    final_output_path = output_path or pptx_path
    
    # 创建临时文件路径，用于处理过程中的临时保存
    temp_pptx = f"{os.path.splitext(pptx_path)[0]}_temp.pptx"
    
    try:
        # 先复制一份原始文件作为中间文件
        import shutil
        shutil.copy2(pptx_path, temp_pptx)
        
        # 记录处理结果
        results = []
        current_slide = slide_number
        
        # 处理每个SVG文件
        for svg_file in svg_files:
            svg_path = os.path.join(svg_dir, svg_file)
            
            # 插入SVG到当前幻灯片
            success = insert_svg_to_pptx(
                pptx_path=temp_pptx,
                svg_path=svg_path,
                slide_number=current_slide,
                output_path=temp_pptx
            )
            
            if success:
                results.append(f"成功: {svg_file} -> 幻灯片 {current_slide}")
                current_slide += 1  # 准备下一张幻灯片
            else:
                results.append(f"失败: {svg_file}")
        
        # 将临时文件重命名为最终输出文件
        if os.path.exists(final_output_path):
            os.remove(final_output_path)
        os.rename(temp_pptx, final_output_path)
        
        # 返回处理结果
        was_created = not os.path.exists(pptx_path) and create_if_not_exists
        creation_msg = "（已自动创建PPTX文件）" if was_created else ""
        using_default_path = not bool(pptx_path) or pptx_path.strip() == ""
        default_path_msg = "（使用自动生成的默认路径）" if using_default_path else ""
        return f"批量处理完成 {creation_msg}{default_path_msg}。处理了 {len(svg_files)} 个SVG文件，结果保存在 {final_output_path}。\n" + "\n".join(results)
    
    except Exception as e:
        # 出现错误时，尝试清理临时文件
        if os.path.exists(temp_pptx):
            try:
                os.remove(temp_pptx)
            except:
                pass
                
        return f"批量处理SVG文件时发生错误: {str(e)}"

# 启动服务器
if __name__ == "__main__":
    mcp.run() 