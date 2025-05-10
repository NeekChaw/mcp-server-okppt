# 检查PPTX文件是否至少有一张幻灯片，如果没有则添加一张
try:
    from pptx import Presentation
    prs = Presentation(pptx_path)
    
    # 检查指定的slide_number是否超出现有幻灯片数量
    if slide_number > len(prs.slides):
        log_error(f"幻灯片编号{slide_number}超出现有幻灯片数量{len(prs.slides)}，将自动添加缺失的幻灯片")
        # 获取空白幻灯片布局
        blank_slide_layout = prs.slide_layouts[6]  # 6是空白幻灯片
        
        # 计算需要添加的幻灯片数量
        slides_to_add = slide_number - len(prs.slides)
        
        # 循环添加所需数量的幻灯片
        for _ in range(slides_to_add):
            prs.slides.add_slide(blank_slide_layout)
            log_error(f"已添加新的空白幻灯片，当前幻灯片数量: {len(prs.slides)}")
        
        # 保存文件
        prs.save(pptx_path)
        # 给文件写入一些时间
        import time
        time.sleep(0.5)
    elif len(prs.slides) == 0:
        log_error(f"PPTX文件 {pptx_path} 没有幻灯片，添加一张空白幻灯片")
        blank_slide_layout = prs.slide_layouts[6]  # 6是空白幻灯片
        slide = prs.slides.add_slide(blank_slide_layout)
        prs.save(pptx_path)
        # 给文件写入一些时间
        import time
        time.sleep(0.5)
except Exception as e:
    error_msg = f"检查或添加幻灯片时出错: {e}"
    log_error(error_msg)
    # 如果是无效的PPTX文件，可能是因为文件损坏或不是PPTX格式
    if "File is not a zip file" in str(e) or "document not found" in str(e) or "Package not found" in str(e):
        log_error(f"PPTX文件 {pptx_path} 似乎不是有效的PowerPoint文件，尝试重新创建")
        try:
            # 确保目录存在
            os.makedirs(os.path.dirname(pptx_path), exist_ok=True)
            
            # 重新创建一个新的PPTX文件
            prs = Presentation()
            prs.slide_width = Inches(16)
            prs.slide_height = Inches(9)
            blank_slide_layout = prs.slide_layouts[6]
            
            # 直接创建足够多的幻灯片
            for i in range(slide_number):
                prs.slides.add_slide(blank_slide_layout)
                
            prs.save(pptx_path)
            log_error(f"已重新创建PPTX文件: {pptx_path}，包含{slide_number}张幻灯片")
            import time
            time.sleep(0.5)
        except Exception as e2:
            error_msg = f"重新创建PPTX文件时出错: {e2}"
            log_error(error_msg)
            log_error(traceback.format_exc())
            return False, "\n".join(error_log)
    else:
        # 其他类型的错误
        log_error(traceback.format_exc())
        return False, "\n".join(error_log) 